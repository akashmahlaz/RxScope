"""Generate RxScope Client Presentation PPT — Clean White Professional Design"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Professional White Color Palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
BODY_TEXT = RGBColor(0x37, 0x41, 0x51)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
PRIMARY = RGBColor(0x25, 0x63, 0xEB)       # Blue
PRIMARY_LIGHT = RGBColor(0xDB, 0xEA, 0xFE) # Light blue bg
SUCCESS = RGBColor(0x05, 0x96, 0x69)       # Green
SUCCESS_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
WARNING = RGBColor(0xD9, 0x77, 0x06)       # Amber
WARNING_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
DANGER = RGBColor(0xDC, 0x26, 0x26)        # Red
DANGER_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)        # Purple

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14, color=BODY_TEXT, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (line_spacing - 1) + 2)
        p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=BODY_TEXT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

def add_line(slide, left, top, width, color=PRIMARY, thickness=3):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(thickness/72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.01))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()

def slide_header(slide, title, subtitle=None):
    add_bg(slide)
    add_line(slide, 1, 0.6, 1.2, PRIMARY, 4)
    add_text(slide, 1, 0.75, 10, 0.7, title, font_size=32, bold=True, color=DARK_TEXT)
    if subtitle:
        add_text(slide, 1, 1.25, 10, 0.4, subtitle, font_size=16, color=MUTED)
    add_divider(slide, 1, 1.65 if subtitle else 1.35, 11.3)

def page_number(slide, num, total=16):
    add_text(slide, 11.5, 7.1, 1.5, 0.3, f"{num} / {total}", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# Blue accent bar at top
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = PRIMARY
top_bar.line.fill.background()

add_line(slide, 1.5, 2.6, 1.5, PRIMARY, 4)
add_text(slide, 1.5, 2.85, 10, 1, "RxScope", font_size=52, bold=True, color=DARK_TEXT)
add_text(slide, 1.5, 3.75, 10, 0.6, "AI-Based Medical Content Identification System", font_size=26, color=PRIMARY)
add_text(slide, 1.5, 4.5, 10, 0.5, "Technical Architecture & Proposal", font_size=18, color=MUTED)
add_divider(slide, 1.5, 5.2, 4)
add_text(slide, 1.5, 5.4, 5, 0.3, "Prepared for RxNetwork", font_size=14, color=BODY_TEXT)
add_text(slide, 1.5, 5.75, 5, 0.3, "April 2026  |  Confidential", font_size=12, color=MUTED)

# ════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "The Problem", "What exists today falls short for HCP-targeted advertising")
page_number(slide, 2)

add_card(slide, 1, 1.9, 5.4, 2.7, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 2.05, 5, 0.35, "Current Industry Gap", font_size=18, bold=True, color=PRIMARY)
add_bullet_list(slide, 1.3, 2.5, 4.9, 2, [
    "Brand safety vendors (IAS, DoubleVerify) classify to avoid — not seek",
    "No solution targets HCP-oriented content specifically",
    "No medical entity resolution (NPI, UMLS, drugs)",
    "No HCP vs consumer content distinction",
], font_size=13, color=DARK_TEXT)

add_card(slide, 6.8, 1.9, 5.5, 2.7, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 7.1, 2.05, 5, 0.35, "What RxNetwork Needs", font_size=18, bold=True, color=SUCCESS)
add_bullet_list(slide, 7.1, 2.5, 5, 2, [
    "AI-driven classification of medically relevant HCP content",
    "Coverage across Scribd, Slideshare & additional platforms",
    "Pharma ad-targeting whitelist in Excel format (CSV/XLSX)",
    "Verified entities: physicians, pharma, medical institutions",
], font_size=13, color=DARK_TEXT)

add_card(slide, 1, 5.0, 11.3, 2.0)
add_text(slide, 1.3, 5.15, 10, 0.35, "RxScope: The Solution", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 1.3, 5.6, 10, 1.2, [
    "Medical taxonomy grounding — MeSH, ICD-10, DSM-5, IAB 3.1 mapping for DSP compatibility",
    "NPI-verified physician identification via federal registry — not guesswork",
    "Document-level AI analysis with ensemble confidence scoring (0–1.0)",
    "Programmatic ad-tech output format — CSV/XLSX whitelist + REST API for real-time DSP integration",
], font_size=13, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 3: HOW IT WORKS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "How RxScope Works", "7-step AI pipeline from URL discovery to whitelist delivery")
page_number(slide, 3)

steps = [
    ("1", "Discover", "Crawl sitemaps &\nURL queues from\ntarget platforms", PRIMARY),
    ("2", "Scrape", "Fetch content via\nheadless browser,\nrobots.txt compliant", PRIMARY),
    ("3", "Extract", "Clean text, metadata\nauthor info from\nHTML / PDF / PPT", PRIMARY),
    ("4", "Classify", "AI classification\nvia Claude + scispaCy\n+ MeSH / IAB mapping", SUCCESS),
    ("5", "Verify", "Entity resolution:\nNPI, RxNorm, FDA\nverification", SUCCESS),
    ("6", "Score", "Ensemble confidence\nscoring + human\nreview for edge cases", WARNING),
    ("7", "Deliver", "CSV/XLSX whitelist\n+ REST API\n+ Admin Dashboard", WARNING),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.55 + i * 1.78
    light = PRIMARY_LIGHT if color == PRIMARY else SUCCESS_LIGHT if color == SUCCESS else WARNING_LIGHT
    add_card(slide, x, 1.85, 1.58, 3.3, light, color)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.55), Inches(2.05), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, x + 0.1, 2.6, 1.4, 0.35, title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.05, 3.05, 1.5, 1.6, desc, font_size=10, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

    if i < 6:
        add_text(slide, x + 1.55, 2.9, 0.3, 0.4, "\u2192", font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

add_card(slide, 1, 5.5, 11.3, 1.5)
add_text(slide, 1.3, 5.6, 10, 0.35, "Key Differentiator", font_size=16, bold=True, color=PRIMARY)
add_multiline(slide, 1.3, 5.95, 10, 0.8, [
    "Multi-agent AI pipeline powered by LangGraph orchestrates specialized agents for each step.",
    "Claude Sonnet handles volume classification. Claude Opus handles edge cases (0.6–0.8 confidence). Human-in-the-loop ensures quality.",
], font_size=12, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 4: TECHNOLOGY STACK
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Technology Stack", "Production-grade, scalable, medical-domain optimized")
page_number(slide, 4)

stack = [
    ("AI & Classification", PRIMARY, PRIMARY_LIGHT, [
        "Claude Sonnet 4 + Opus 4 (Anthropic)",
        "LangGraph — agent orchestration (MIT, $0)",
        "LangSmith — monitoring & observability",
        "scispaCy + MedSpaCy — medical NLP",
        "BiomedBERT — biomedical embeddings",
    ]),
    ("Data & Storage", SUCCESS, SUCCESS_LIGHT, [
        "PostgreSQL 16 + pgvector",
        "Managed cloud database hosting",
        "Redis — caching & job queues",
        "Cloudflare R2 — object storage",
    ]),
    ("Infrastructure", WARNING, WARNING_LIGHT, [
        "Lightpanda — headless browser scraping",
        "Cloudflare Workers — edge API",
        "Cloudflare Pages — dashboard hosting",
        "Rotating residential proxies",
    ]),
    ("Medical APIs (All Free)", ACCENT, RGBColor(0xED, 0xE9, 0xFE), [
        "NPI Registry — physician verification",
        "RxNorm — drug name resolution",
        "FDA NDC — drug identification",
        "UMLS / MeSH — medical ontologies",
    ]),
]

for i, (title, color, bg, items) in enumerate(stack):
    x = 0.7 + i * 3.1
    add_card(slide, x, 1.85, 2.9, 3.5, bg, color)
    add_text(slide, x + 0.15, 1.95, 2.6, 0.35, title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_line(slide, x + 0.4, 2.35, 2.1, color, 2)
    add_bullet_list(slide, x + 0.1, 2.5, 2.7, 2.6, items, font_size=11, color=BODY_TEXT)

add_card(slide, 0.7, 5.6, 11.9, 1.4)
add_text(slide, 1, 5.7, 5, 0.3, "LangGraph = the ENGINE", font_size=14, bold=True, color=PRIMARY)
add_text(slide, 1, 6.0, 5.3, 0.7, "Open-source framework that builds the AI pipeline.\nOrchestrates agents: scrape \u2192 classify \u2192 score \u2192 store.\nCost: $0 (MIT license)", font_size=11, color=BODY_TEXT)
add_text(slide, 6.5, 5.7, 5.5, 0.3, "LangSmith = the MONITOR", font_size=14, bold=True, color=SUCCESS)
add_text(slide, 6.5, 6.0, 5.3, 0.7, "SaaS platform that watches the pipeline in production.\nTraces every decision, measures cost & accuracy.\nFree tier available; scales with usage.", font_size=11, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 5: DATABASE DECISION
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Database Architecture", "PostgreSQL + pgvector — unified relational + vector database")
page_number(slide, 5)

add_card(slide, 1, 1.85, 5.5, 3.0, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 1.3, 1.95, 5, 0.35, "\u2713  PostgreSQL + pgvector (Selected)", font_size=17, bold=True, color=SUCCESS)
add_bullet_list(slide, 1.3, 2.4, 5, 2.3, [
    "Native vector search — single DB for everything",
    "Full ACID compliance — guaranteed data consistency",
    "Foreign keys for medical entity relationships",
    "SQL \u2192 CSV/XLSX export (client requirement)",
    "LangGraph first-class checkpointer support",
    "Connection pooling via Hyperdrive (Cloudflare)",
], font_size=12, color=DARK_TEXT)

add_card(slide, 6.8, 1.85, 5.5, 3.0, DANGER_LIGHT, DANGER)
add_text(slide, 7.1, 1.95, 5, 0.35, "\u2717  MongoDB (Rejected)", font_size=17, bold=True, color=DANGER)
add_bullet_list(slide, 7.1, 2.4, 5, 2.3, [
    "Separate vector search service needed ($$$)",
    "No JOINs — denormalized data management",
    "Limited multi-document transactions",
    "Complex aggregation pipelines for CSV export",
    "No official LangGraph checkpointer",
    "Wrong paradigm for relational medical data",
], font_size=12, color=DARK_TEXT)

add_card(slide, 1, 5.15, 11.3, 2.0)
add_text(slide, 1.3, 5.25, 10, 0.35, "Cloud Database Hosting", font_size=17, bold=True, color=PRIMARY)
add_multiline(slide, 1.3, 5.65, 10, 1.3, [
    "Managed PostgreSQL with pgvector extension  \u2022  Automatic backups & connection pooling",
    "Scalable compute from shared instances ($25/mo) to dedicated 16XL ($3,730/mo)",
    "8GB disk included in base plan  \u2022  Scales seamlessly as data grows  \u2022  Zero-downtime upgrades",
], font_size=12, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 6: WHAT YOU GET
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "What RxNetwork Gets", "Complete end-to-end medical content classification system")
page_number(slide, 6)

deliverables = [
    ("AI Classification", "Automated identification of HCP-oriented\ncontent across multiple platforms", PRIMARY, PRIMARY_LIGHT),
    ("Excel Whitelist", "Deduplicated URL list in CSV/XLSX\nwith full metadata & confidence scores", PRIMARY, PRIMARY_LIGHT),
    ("Admin Dashboard", "Real-time monitoring, search, export,\nand human-in-the-loop review panel", SUCCESS, SUCCESS_LIGHT),
    ("REST API", "Programmatic whitelist access for DSP\nintegration (OpenRTB compatible)", SUCCESS, SUCCESS_LIGHT),
    ("Entity Verification", "NPI-verified physicians, FDA-validated\ndrugs, verified medical institutions", WARNING, WARNING_LIGHT),
    ("Multi-Taxonomy", "IAB 3.1, MeSH, ICD-10, DSM-5, RxNorm\n— all mapped & cross-referenced", WARNING, WARNING_LIGHT),
    ("Re-validation", "Scheduled re-scraping ensures URLs\nremain current and accurately classified", ACCENT, RGBColor(0xED, 0xE9, 0xFE)),
    ("Audit Trail", "Every classification decision is\nfully traceable and explainable", ACCENT, RGBColor(0xED, 0xE9, 0xFE)),
]

for i, (title, desc, color, bg) in enumerate(deliverables):
    col = i % 4
    row = i // 4
    x = 0.7 + col * 3.1
    y = 1.85 + row * 2.55
    add_card(slide, x, y, 2.9, 2.2, bg, color)
    add_text(slide, x + 0.15, y + 0.2, 2.6, 0.3, title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_line(slide, x + 0.5, y + 0.55, 1.9, color, 2)
    add_text(slide, x + 0.15, y + 0.7, 2.6, 1.2, desc, font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 7: ADMIN DASHBOARD
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Admin Dashboard", "Web-based control center for monitoring, review, and export")
page_number(slide, 7)

pages = [
    ("Overview", "Total URLs processed, whitelist count,\nconfidence distribution, queue status,\nsystem health metrics", PRIMARY),
    ("Whitelist Explorer", "Searchable and filterable table,\nexport to CSV/XLSX, confidence and\ncategory filters, inline preview", PRIMARY),
    ("HITL Review", "Low-confidence review queue,\none-click approve/reject,\noverride with audit reason", SUCCESS),
    ("Entity Browser", "Browse by: pharmaceutical companies,\nphysicians (NPI-verified), medical\nschools, biotech, institutions", SUCCESS),
    ("Analytics", "Accuracy over time, API cost tracking,\nprocessing speed metrics, content\nchange rates, exportable reports", WARNING),
    ("Settings", "Confidence thresholds, platform\nmanagement, re-scrape scheduling,\nAPI keys, user roles", WARNING),
]

for i, (title, desc, color) in enumerate(pages):
    col = i % 3
    row = i // 3
    x = 0.7 + col * 4.05
    y = 1.85 + row * 2.55
    bg = PRIMARY_LIGHT if color == PRIMARY else SUCCESS_LIGHT if color == SUCCESS else WARNING_LIGHT
    add_card(slide, x, y, 3.75, 2.2, bg, color)
    add_text(slide, x + 0.2, y + 0.15, 3.3, 0.35, title, font_size=16, bold=True, color=color)
    add_line(slide, x + 0.2, y + 0.5, 2, color, 2)
    add_text(slide, x + 0.2, y + 0.65, 3.3, 1.4, desc, font_size=12, color=BODY_TEXT)

add_text(slide, 1, 7.0, 11, 0.3, "Built with Next.js  \u2022  Hosted on Cloudflare Pages ($0)  \u2022  Real-time updates  \u2022  Secure authentication", font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 8: OUTPUT FORMAT
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Output: Whitelist Format", "Excel-compatible CSV/XLSX per RxNetwork specification")
page_number(slide, 8)

columns = [
    ("url", "Deduplicated URL validated as HCP content"),
    ("source_domain", "Platform domain (e.g., scribd.com, slideshare.net)"),
    ("content_type", "PDF, PPT, white paper, clinical summary, research poster, educational material"),
    ("detected_medical_categories", "MeSH codes + IAB categories (semicolon-separated)"),
    ("detected_entity_associations", "Organization names, physician names, drug names"),
    ("confidence_score", "Overall classification confidence (0.00 – 1.00)"),
]

add_card(slide, 1, 1.85, 11.3, 3.8, CARD_BG, BORDER)
add_text(slide, 1.3, 1.95, 10, 0.35, "Required Output Columns", font_size=17, bold=True, color=PRIMARY)
add_line(slide, 1.3, 2.3, 3, PRIMARY, 2)

# Table header
add_text(slide, 1.5, 2.5, 3, 0.3, "Column", font_size=12, bold=True, color=MUTED)
add_text(slide, 4.8, 2.5, 7, 0.3, "Description", font_size=12, bold=True, color=MUTED)
add_divider(slide, 1.5, 2.75, 10.3)

for i, (col, desc) in enumerate(columns):
    y = 2.85 + i * 0.42
    add_text(slide, 1.5, y, 3, 0.35, col, font_size=12, bold=True, color=PRIMARY)
    add_text(slide, 4.8, y, 7, 0.35, desc, font_size=12, color=BODY_TEXT)
    if i < 5:
        add_divider(slide, 1.5, y + 0.35, 10.3)

add_card(slide, 1, 5.95, 11.3, 1.2, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 6.05, 10, 0.3, "Extended Columns (Recommended)", font_size=15, bold=True, color=PRIMARY)
add_text(slide, 1.3, 6.4, 10, 0.5, "audience_type  \u2022  source_type  \u2022  iab_category_codes  \u2022  icd10_codes  \u2022  dsm5_categories  \u2022  fda_drug_names  \u2022  mesh_codes  \u2022  last_validated_at", font_size=12, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 9: ENTITY CATEGORIES
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Target Entity Categories", "8 categories per RxNetwork specification — all verified against authoritative sources")
page_number(slide, 9)

entities = [
    ("Pharmaceutical\nManufacturers", "FDA-verified pharma\ncompany identification", PRIMARY, PRIMARY_LIGHT),
    ("Medical Schools &\nAcademic Centers", "Accredited medical education\ninstitutions", PRIMARY, PRIMARY_LIGHT),
    ("Medical Trade\nAssociations", "AMA, specialty societies,\nprofessional organizations", PRIMARY, PRIMARY_LIGHT),
    ("Recognized\nPhysicians", "NPI Registry verified,\nlicense-confirmed MDs/DOs", PRIMARY, PRIMARY_LIGHT),
    ("Healthcare\nExecutives", "C-suite and leadership\nin healthcare organizations", SUCCESS, SUCCESS_LIGHT),
    ("Biotech\nCompanies", "Biotechnology firms,\nresearch organizations", SUCCESS, SUCCESS_LIGHT),
    ("Medical Device\nManufacturers", "FDA-listed medical device\ncompanies", SUCCESS, SUCCESS_LIGHT),
    ("Prominent Medical\nInstitutions", "Major hospitals, research\ncenters, health systems", SUCCESS, SUCCESS_LIGHT),
]

for i, (title, desc, color, bg) in enumerate(entities):
    col = i % 4
    row = i // 4
    x = 0.7 + col * 3.1
    y = 1.85 + row * 2.5
    add_card(slide, x, y, 2.9, 2.15, bg, color)
    add_text(slide, x + 0.15, y + 0.2, 2.6, 0.5, title, font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_line(slide, x + 0.5, y + 0.75, 1.9, color, 1.5)
    add_text(slide, x + 0.15, y + 0.9, 2.6, 1, desc, font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 10: CLOUDFLARE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Cloud Infrastructure", "Selective use of Cloudflare services — verified against official documentation")
page_number(slide, 10)

add_text(slide, 1, 1.85, 5, 0.35, "Services We Use", font_size=18, bold=True, color=SUCCESS)

use_items = [
    ("Workers", "Whitelist REST API at global edge", "$5/mo for 10M requests"),
    ("R2 Storage", "Document & content archival", "$0.015/GB, zero egress fees"),
    ("Hyperdrive", "PostgreSQL connection pooler", "Included on paid plan"),
    ("Pages", "Admin dashboard hosting", "$0 — free static hosting"),
]

for i, (name, desc, cost) in enumerate(use_items):
    y = 2.3 + i * 0.7
    add_card(slide, 1, y, 5.5, 0.6, SUCCESS_LIGHT, SUCCESS)
    add_text(slide, 1.2, y + 0.1, 1.5, 0.3, name, font_size=13, bold=True, color=SUCCESS)
    add_text(slide, 2.8, y + 0.1, 2, 0.3, desc, font_size=11, color=BODY_TEXT)
    add_text(slide, 4.8, y + 0.1, 1.6, 0.3, cost, font_size=10, bold=True, color=MUTED, alignment=PP_ALIGN.RIGHT)

add_text(slide, 7, 1.85, 5, 0.35, "Services We Don't Use (and Why)", font_size=18, bold=True, color=DANGER)

dont_items = [
    ("Browser Rendering", "Our scraping engine is 50x cheaper"),
    ("Workers AI", "No biomedical classification models available"),
    ("D1 Database", "SQLite-based — cannot replace PostgreSQL"),
    ("Vectorize", "Cost-prohibitive at scale ($19K/mo at 50M URLs)"),
]

for i, (name, reason) in enumerate(dont_items):
    y = 2.3 + i * 0.7
    add_card(slide, 6.8, y, 5.5, 0.6, DANGER_LIGHT, DANGER)
    add_text(slide, 7.0, y + 0.1, 1.8, 0.3, name, font_size=13, bold=True, color=DANGER)
    add_text(slide, 8.8, y + 0.1, 3.3, 0.3, reason, font_size=11, color=BODY_TEXT)

add_card(slide, 1, 5.2, 11.3, 1.0)
add_text(slide, 1.3, 5.3, 10, 0.3, "All pricing verified against official Cloudflare documentation (April 2026)", font_size=12, bold=True, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, 1.3, 5.6, 10, 0.3, "Workers: $5/mo + $0.30/M requests   \u2022   R2: $0.015/GB + $0 egress   \u2022   Hyperdrive: Free   \u2022   Pages: Free", font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 11: DEVELOPMENT INVESTMENT
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Development Investment", "One-time development cost — milestone-based delivery")
page_number(slide, 11)

phases = [
    ("Phase 0", "Infrastructure & Setup", "DB schema, evaluation dataset, CI/CD", "1–2 weeks"),
    ("Phase 1", "Scraping Pipeline", "Lightpanda, content extraction, proxy setup", "2–3 weeks"),
    ("Phase 2", "Classification Engine", "LangGraph + Claude + scispaCy + NLP pipeline", "3–4 weeks"),
    ("Phase 3", "Entity Resolution", "NPI, RxNorm, FDA, UMLS integration", "2–3 weeks"),
    ("Phase 4", "Dashboard & Export", "Admin panel, CSV/XLSX export, HITL review", "1–2 weeks"),
    ("Phase 5", "Tuning & QA", "Accuracy tuning, testing, security audit", "2–3 weeks"),
    ("Phase 6", "Deployment & Docs", "Production deployment, documentation, handoff", "1–2 weeks"),
]

# Column headers
add_text(slide, 1.3, 1.85, 1.5, 0.3, "Phase", font_size=12, bold=True, color=MUTED)
add_text(slide, 2.8, 1.85, 2.5, 0.3, "Deliverable", font_size=12, bold=True, color=MUTED)
add_text(slide, 5.8, 1.85, 3.5, 0.3, "Scope", font_size=12, bold=True, color=MUTED)
add_text(slide, 10, 1.85, 2, 0.3, "Duration", font_size=12, bold=True, color=MUTED, alignment=PP_ALIGN.RIGHT)
add_divider(slide, 1.3, 2.15, 10.7)

for i, (phase, deliverable, scope, duration) in enumerate(phases):
    y = 2.3 + i * 0.52
    bg_color = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, 1, y - 0.05, 11.3, 0.48, bg_color, bg_color)
    add_text(slide, 1.3, y, 1.3, 0.3, phase, font_size=12, bold=True, color=PRIMARY)
    add_text(slide, 2.8, y, 2.5, 0.3, deliverable, font_size=12, bold=True, color=DARK_TEXT)
    add_text(slide, 5.8, y, 3.8, 0.3, scope, font_size=11, color=BODY_TEXT)
    add_text(slide, 10, y, 2, 0.3, duration, font_size=12, color=MUTED, alignment=PP_ALIGN.RIGHT)

# Total row
add_card(slide, 1, 5.95, 11.3, 0.55, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 6.05, 3, 0.3, "Total Development", font_size=14, bold=True, color=PRIMARY)
add_text(slide, 10, 6.05, 2, 0.3, "12–19 weeks", font_size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.RIGHT)

add_card(slide, 1, 6.7, 11.3, 0.6)
add_text(slide, 1.3, 6.75, 10, 0.4, "Recommended: Milestone-based payments — 20% upfront, remainder tied to deliverable milestones (M1–M5)", font_size=13, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 12: MONTHLY INFRASTRUCTURE COSTS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Monthly Infrastructure Costs", "Recurring costs scale with URL volume — client pays directly to service providers")
page_number(slide, 12)

# Column headers
for ci, (label, x, color, bg) in enumerate([
    ("Small — 100K URLs", 1.3, PRIMARY, PRIMARY_LIGHT),
    ("Medium — 1M URLs", 5.3, SUCCESS, SUCCESS_LIGHT),
    ("Large — 50M URLs", 9.3, WARNING, WARNING_LIGHT),
]):
    add_card(slide, x - 0.3, 1.85, 3.8, 0.5, bg, color)
    add_text(slide, x, 1.9, 3.2, 0.35, label, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)

services = [
    ("Database (PostgreSQL)", "$25", "$40", "$275–400"),
    ("AI Classification (Claude)", "$50–100", "$500–1,000", "$10.5K–25K"),
    ("Scraping Infrastructure", "$20", "$80", "$200–500"),
    ("Cache & Queues (Redis)", "$0", "$10", "$50–100"),
    ("Edge API (Cloudflare)", "$5", "$5", "$5–45"),
    ("Object Storage", "$0", "$5", "$15–75"),
    ("Monitoring (LangSmith)", "$0", "$164", "$300–500"),
    ("Rotating Proxies", "$50–100", "$100–200", "$500–1K"),
]

add_text(slide, 1.3, 2.5, 2.5, 0.3, "Service", font_size=11, bold=True, color=MUTED)
add_divider(slide, 1, 2.75, 11.5)

for i, (service, small, med, large) in enumerate(services):
    y = 2.85 + i * 0.4
    bg_color = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, 1, y - 0.05, 11.5, 0.38, bg_color, bg_color)
    add_text(slide, 1.3, y, 2.8, 0.3, service, font_size=11, color=BODY_TEXT)
    add_text(slide, 4.2, y, 1.4, 0.3, small, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text(slide, 8.2, y, 1.4, 0.3, med, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text(slide, 12.0, y, 1.2, 0.3, large, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

# Totals
y_total = 2.85 + len(services) * 0.4 + 0.15
add_line(slide, 1, y_total - 0.05, 11.5, BORDER, 2)
add_text(slide, 1.3, y_total, 2.5, 0.3, "Monthly Total", font_size=13, bold=True, color=DARK_TEXT)
add_text(slide, 3.5, y_total, 2, 0.3, "$170–300", font_size=13, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_text(slide, 7.5, y_total, 2, 0.3, "$1,004–1,704", font_size=13, bold=True, color=SUCCESS, alignment=PP_ALIGN.CENTER)
add_text(slide, 11.2, y_total, 2, 0.3, "$12K–28K", font_size=13, bold=True, color=WARNING, alignment=PP_ALIGN.RIGHT)

add_card(slide, 1, 6.6, 11.3, 0.7, WARNING_LIGHT, WARNING)
add_text(slide, 1.3, 6.65, 10, 0.5, "Note: Claude API is the dominant cost at scale. Pre-filtering heuristics and batch processing can reduce AI costs by 50–70%.", font_size=12, color=WARNING, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 13: OTHER EXPENSES & SETUP COSTS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Additional Expenses", "One-time setup costs and optional services")
page_number(slide, 13)

add_text(slide, 1, 1.85, 5, 0.35, "One-Time Setup Costs", font_size=18, bold=True, color=PRIMARY)

setup_items = [
    ("Claude API initial deposit", "$50–100", "Pre-paid credits for AI classification"),
    ("Proxy service setup", "$50–100", "Initial deposit for residential proxy plan"),
    ("Domain registration", "$10–15/year", "Custom domain for dashboard & API"),
    ("SSL Certificate", "$0", "Provided free by Cloudflare"),
    ("UMLS License", "$0", "Free federal license (1–3 day approval)"),
]

for i, (item, cost, note) in enumerate(setup_items):
    y = 2.3 + i * 0.52
    bg_color = CARD_BG if i % 2 == 0 else WHITE
    add_card(slide, 1, y - 0.05, 5.5, 0.48, bg_color, BORDER)
    add_text(slide, 1.2, y, 2.5, 0.3, item, font_size=12, color=DARK_TEXT)
    add_text(slide, 3.7, y, 1, 0.3, cost, font_size=12, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_text(slide, 4.8, y, 1.6, 0.3, note, font_size=10, color=MUTED)

add_card(slide, 1, 5.0, 5.5, 0.5, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 5.05, 3, 0.3, "Total Setup:", font_size=13, bold=True, color=PRIMARY)
add_text(slide, 3.7, 5.05, 1, 0.3, "$110–215", font_size=13, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

add_text(slide, 7, 1.85, 5, 0.35, "Optional Services", font_size=18, bold=True, color=WARNING)

optional_items = [
    ("Legal review of scraping approach", "$500–2,000", "Recommended before Scribd/SlideShare"),
    ("Evaluation dataset labeling", "$500–1,000", "If outsourced (included if in-house)"),
    ("Post-launch maintenance", "$1,000–2,000/mo", "Monitoring, bug fixes, re-tuning"),
    ("Scale optimization consulting", "Custom", "Performance tuning at 10M+ URLs"),
]

for i, (item, cost, note) in enumerate(optional_items):
    y = 2.3 + i * 0.6
    add_card(slide, 6.8, y - 0.05, 5.5, 0.52, WARNING_LIGHT, WARNING)
    add_text(slide, 7.0, y, 3, 0.3, item, font_size=12, color=DARK_TEXT)
    add_text(slide, 10.0, y, 1.3, 0.3, cost, font_size=11, bold=True, color=WARNING, alignment=PP_ALIGN.RIGHT)
    add_text(slide, 7.0, y + 0.22, 4, 0.25, note, font_size=10, color=MUTED)

# ════════════════════════════════════════════
# SLIDE 14: TOTAL INVESTMENT SUMMARY
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Investment Summary", "Complete cost overview for the RxScope system")
page_number(slide, 14)

add_card(slide, 1, 1.85, 3.6, 2.5, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.2, 2.0, 3.2, 0.35, "Development", font_size=13, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_text(slide, 1.2, 2.35, 3.2, 0.5, "$18,500 – $30,500", font_size=22, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text(slide, 1.2, 2.85, 3.2, 0.3, "One-time", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, 1.2, 3.15, 3.2, 0.5, "12–19 weeks, milestone-based\npayments tied to deliverables", font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

add_card(slide, 4.9, 1.85, 3.6, 2.5, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 5.1, 2.0, 3.2, 0.35, "Infrastructure", font_size=13, bold=True, color=SUCCESS, alignment=PP_ALIGN.CENTER)
add_text(slide, 5.1, 2.35, 3.2, 0.5, "$170 – $1,704", font_size=22, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text(slide, 5.1, 2.85, 3.2, 0.3, "Per month (100K–1M URLs)", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, 5.1, 3.15, 3.2, 0.5, "Paid directly to cloud providers.\nScales with URL volume.", font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

add_card(slide, 8.8, 1.85, 3.6, 2.5, WARNING_LIGHT, WARNING)
add_text(slide, 9.0, 2.0, 3.2, 0.35, "Setup & Extras", font_size=13, bold=True, color=WARNING, alignment=PP_ALIGN.CENTER)
add_text(slide, 9.0, 2.35, 3.2, 0.5, "$110 – $215", font_size=22, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text(slide, 9.0, 2.85, 3.2, 0.3, "One-time", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, 9.0, 3.15, 3.2, 0.5, "API deposits, domain,\nproxy initial setup", font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# Payment structure
add_card(slide, 1, 4.7, 11.3, 2.5)
add_text(slide, 1.3, 4.8, 10, 0.35, "Recommended Payment Structure (Milestone-Based)", font_size=16, bold=True, color=PRIMARY)

milestones_pay = [
    ("20%", "Project Kickoff", "Infrastructure setup, DB schema, evaluation dataset begin", "Week 0"),
    ("20%", "Proof of Concept", "Working classifier on 3,000 PubMed URLs + accuracy report", "Week 4–5"),
    ("25%", "First Whitelist", "100K URL whitelist delivered in CSV/XLSX format", "Week 11"),
    ("20%", "Dashboard Live", "Admin panel with search, export, and HITL review", "Week 14"),
    ("15%", "Production Launch", "Full system running at target scale with monitoring", "Week 18–22"),
]

add_text(slide, 1.5, 5.2, 0.8, 0.3, "%", font_size=10, bold=True, color=MUTED)
add_text(slide, 2.3, 5.2, 2.2, 0.3, "Milestone", font_size=10, bold=True, color=MUTED)
add_text(slide, 4.8, 5.2, 5, 0.3, "Deliverable", font_size=10, bold=True, color=MUTED)
add_text(slide, 10.5, 5.2, 1.5, 0.3, "When", font_size=10, bold=True, color=MUTED, alignment=PP_ALIGN.RIGHT)
add_divider(slide, 1.5, 5.45, 10.5)

for i, (pct, name, desc, when) in enumerate(milestones_pay):
    y = 5.55 + i * 0.33
    add_text(slide, 1.5, y, 0.6, 0.3, pct, font_size=11, bold=True, color=PRIMARY)
    add_text(slide, 2.3, y, 2.2, 0.3, name, font_size=11, bold=True, color=DARK_TEXT)
    add_text(slide, 4.8, y, 5.2, 0.3, desc, font_size=10, color=BODY_TEXT)
    add_text(slide, 10.5, y, 1.5, 0.3, when, font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════
# SLIDE 15: IMPORTANT NOTES
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Important Considerations", "Technical realities to be aware of before project kick-off")
page_number(slide, 15)

notes = [
    ("Scribd / SlideShare Access", "Scribd has no public API since 2017. Content access is limited to metadata scraping. We recommend starting with legally clean sources (PubMed, open-access journals) in Phase 1, adding Scribd/SlideShare in Phase 2.", WARNING, WARNING_LIGHT),
    ("Classification Accuracy", "Expect 85–90% accuracy initially, improving to 95%+ after human-in-the-loop feedback cycles. This is standard for AI classification systems — there is a 2–4 week tuning period.", PRIMARY, PRIMARY_LIGHT),
    ("Cost Scales with Volume", "At 100K URLs: ~$200/month infrastructure. At 50M URLs: ~$15K–20K/month (dominated by Claude API). Pre-filtering heuristics can reduce AI costs by 50–70%.", SUCCESS, SUCCESS_LIGHT),
    ("robots.txt Compliance", "We respect robots.txt on all platforms. Some sites may limit scraping rate. This is a legal requirement and protects the project long-term.", PRIMARY, PRIMARY_LIGHT),
    ("Ongoing Maintenance", "Medical content changes constantly — new drugs, guidelines, HCPs. Whitelists need periodic re-scraping and re-classification. Post-launch maintenance recommended.", WARNING, WARNING_LIGHT),
    ("Separate Billing", "Development is a one-time cost. Infrastructure is monthly recurring, billed directly by service providers (Anthropic, Cloudflare, etc.) to your accounts.", SUCCESS, SUCCESS_LIGHT),
]

for i, (title, desc, color, bg) in enumerate(notes):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 6.2
    y = 1.85 + row * 1.75
    add_card(slide, x, y, 5.9, 1.5, bg, color)
    add_text(slide, x + 0.2, y + 0.1, 5.5, 0.3, title, font_size=14, bold=True, color=color)
    add_text(slide, x + 0.2, y + 0.45, 5.5, 0.9, desc, font_size=11, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 16: NEXT STEPS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Next Steps", "Questions to finalize scope, and immediate action items")
page_number(slide, 16)

add_card(slide, 1, 1.85, 5.5, 4.2, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 1.95, 5, 0.35, "Questions We Need From You", font_size=17, bold=True, color=PRIMARY)
add_line(slide, 1.3, 2.3, 2, PRIMARY, 2)
add_bullet_list(slide, 1.3, 2.5, 5, 3.5, [
    "How many URLs in the first batch?",
    "One-time batch or continuously running system?",
    "Dashboard required, or CSV exports only?",
    "Who provides the seed URL list?",
    "Expected turnaround time?",
    "Which platforms are mandatory for Phase 1?",
    "Specific DSP integration needed?",
    "Budget range for development?",
], font_size=12, color=DARK_TEXT)

add_card(slide, 6.8, 1.85, 5.5, 4.2, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 7.1, 1.95, 5, 0.35, "Once Confirmed — We Begin", font_size=17, bold=True, color=SUCCESS)
add_line(slide, 7.1, 2.3, 2, SUCCESS, 2)
add_bullet_list(slide, 7.1, 2.5, 5, 3.5, [
    "Confirm URL volume to finalize cost quote",
    "Agree on milestone-based payment schedule",
    "Confirm Phase 1 platform scope",
    "Sign-off on technology stack",
    "Set up shared cloud accounts",
    "Kick off Phase 0: Foundation",
    "Schedule weekly progress demos",
    "Begin 3,000-URL evaluation dataset",
], font_size=12, color=DARK_TEXT)

add_card(slide, 1, 6.35, 11.3, 0.8, CARD_BG, PRIMARY)
add_text(slide, 1.3, 6.4, 10, 0.3, "Ready to Build", font_size=18, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_text(slide, 1.3, 6.7, 10, 0.3, "Architecture is researched & documented  \u2022  Technology stack selected  \u2022  Cost models built  \u2022  Awaiting your confirmation to start", font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════
output_path = r"c:\Users\akash\work\rxscope\RxScope_Client_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
